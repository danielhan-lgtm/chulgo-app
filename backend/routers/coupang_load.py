# -*- coding: utf-8 -*-
"""쿠팡 팔레트 적재리스트(PPT) 생성기 (API 0원)

거래명세서(PDF) + 부착리스트(PDF)를 업로드하면 내용을 파싱해
쿠팡 제출용 '적재리스트' PPT를 자동 생성한다.

- 거래명세서  → 업체명/납품센터/입고예정일 + SKU 표(상품번호·상품명·확정수량·유통기한)
- 부착리스트  → 입고요청서번호(밀크런 번호) + 팔레트 번호
- BOX 수량    → 거래명세서에 없는 값(입수 의존)이라 화면에서 사용자가 입력/수정한다.

PPT는 templates/쿠팡_적재리스트_양식.pptx 를 베이스로 동적 항목만 치환한다.
외부 API를 호출하지 않는다.
"""
import copy
import io
import os
import re
import sys

import pdfplumber
from fastapi import APIRouter, File, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
from pydantic import BaseModel

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import load_settings as LS

router = APIRouter()

TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "templates", "쿠팡_적재리스트_양식.pptx"
)


# ── 정규화 유틸 ────────────────────────────────────────────────────────────
def _norm_center(name: str) -> str:
    """'창원4(06)' → '창원4'"""
    if not name:
        return ""
    s = re.sub(r"\s*\(.*?\)\s*$", "", str(name).strip())
    return re.sub(r"\s+", "", s)


def _norm_product(name: str) -> str:
    return re.sub(r"\s+", " ", str(name or "")).strip()


def _to_int(v) -> int:
    try:
        return int(re.sub(r"[^\d\-]", "", str(v)) or 0)
    except Exception:
        return 0


# ── BOX 수량(입수) 자동 추정 ─────────────────────────────────────────────────
# 거래명세서에는 입수 정보가 없어 저장된 설정(load_settings.json)으로 추정한다.
# 규칙(상품번호/키워드별 1박스당 수량) → 번들 표기 → 기본 입수 순. 화면에서 수정 가능.


# ── 거래명세서 파서 ────────────────────────────────────────────────────────
def _parse_거래명세서_국문(data: bytes, settings: dict) -> dict:
    """2026-06 이후 쿠팡 국문 양식 — '거래 명세서 / 쿠팡 제출용' (물류동봉문서).

    페이지가 (쿠팡 제출용/업체 보관용) 쌍으로 반복 → 발주번호 기준 dedup.
    발주정보 표에서 업체명·납품 센터·도착예정일·팔레트수량,
    상품 표에서 상품번호·상품명·확정수량·소비기한을 뽑는다.
    """
    supplier = center = date = None
    pallet_total = 0
    seen_po: set[str] = set()
    ordered_rows: list[dict] = []

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pg in pdf.pages:
            txt = pg.extract_text() or ""
            mpo = re.search(r"발주번호\s*(\d+)", txt)
            po = mpo.group(1) if mpo else None
            if po and po in seen_po:
                continue  # 업체 보관용(중복) 페이지
            if po:
                seen_po.add(po)

            kv: dict[str, str] = {}
            item_tbl = None
            for tbl in pg.extract_tables() or []:
                if not tbl:
                    continue
                header = [(c or "").strip() for c in tbl[0]]
                if (any(("상품번호" in h or "상품코드" in h) for h in header)
                        and any("발주수량" in h for h in header)):
                    item_tbl = tbl
                    continue
                for row in tbl:
                    if row and len(row) >= 2 and row[0]:
                        kv[str(row[0]).strip()] = str(row[1] or "").strip()

            if supplier is None:
                raw_sup = kv.get("업체명") or kv.get("거래처명") or ""
                if raw_sup:
                    supplier = re.sub(r"^\[[^\]]+\]\s*", "", raw_sup)  # '[A0125..] 상호' → '상호'
            center = center or kv.get("납품 센터") or kv.get("납품센터") or kv.get("FC명")
            if not date:
                # '물류센터 도착예정일 20260717' 또는 '입고예정일 2026-06-30'
                raw = re.sub(r"\D", "", kv.get("물류센터 도착예정일") or kv.get("입고예정일") or "")
                if len(raw) >= 8:
                    date = f"{raw[:4]}-{raw[4:6]}-{raw[6:8]}"
            if not pallet_total:
                pallet_total = _to_int(kv.get("팔레트수량", 0))

            if not item_tbl:
                continue
            header = [(c or "").replace("\n", "") for c in item_tbl[0]]
            i_no = next((i for i, h in enumerate(header) if h.strip() == "No"), 0)
            i_sku = next((i for i, h in enumerate(header)
                          if "상품번호" in h or "상품코드" in h), 1)
            i_name = next((i for i, h in enumerate(header) if "상품명" in h), i_sku)
            i_qty = next((i for i, h in enumerate(header) if "발주수량" in h), 3)
            i_recv = next((i for i, h in enumerate(header) if "확정수량" in h), None)
            # 유통/소비기한 컬럼 — '...관리'(Y/N 플래그) 컬럼은 제외
            i_exp = next((i for i, h in enumerate(header)
                          if ("소비기한" in h or "유통" in h) and "관리" not in h), None)
            for row in item_tbl[1:]:
                if not row or len(row) <= max(i_sku, i_qty):
                    continue
                c0 = str(row[i_no] or "").strip()
                if not c0.isdigit():
                    continue  # '합계' 등
                cell_lines = [p.strip() for p in str(row[i_sku] or "").split("\n") if p.strip()]
                if i_sku == i_name:
                    # 결합 셀: '상품코드\n상품명...\nBarcode\nBOX바코드'
                    if not cell_lines or not cell_lines[0].isdigit():
                        continue
                    sku = cell_lines[0]
                    name_src = cell_lines[1:]
                else:
                    sku = cell_lines[0] if cell_lines else ""
                    name_src = [p.strip() for p in str(row[i_name] or "").split("\n") if p.strip()]
                # 이름 줄들 — 바코드/‘-’ 줄 전까지 이어붙임
                name_lines = []
                for p in name_src:
                    if p == "-" or re.fullmatch(r"[A-Z]?\d{9,}", p):
                        break
                    name_lines.append(p)
                name = " ".join(name_lines)
                # 확정수량 컬럼이 있으면 그 값 사용 (0 = 미확정/취소 → 제외), 없으면 발주수량
                if i_recv is not None:
                    qty = _to_int(row[i_recv] if i_recv < len(row) else 0)
                else:
                    qty = _to_int(row[i_qty])
                if qty <= 0:
                    continue
                expire = ""
                exp_cell = str(row[i_exp] or "") if (i_exp is not None and i_exp < len(row)) else ""
                found = re.findall(r"\d{4}-\d{2}-\d{2}|\d{8}", exp_cell)
                if found:
                    d = found[-1]  # (제조일자 ~ 소비기한) 중 마지막 = 소비기한
                    expire = d if "-" in d else f"{d[:4]}-{d[4:6]}-{d[6:]}"
                ordered_rows.append({
                    "sku": sku, "name": _norm_product(name), "qty": qty, "expire": expire,
                })

    merged: dict[str, dict] = {}
    order: list[str] = []
    for r in ordered_rows:
        if r["sku"] in merged:
            merged[r["sku"]]["qty"] += r["qty"]
            if not merged[r["sku"]]["expire"]:
                merged[r["sku"]]["expire"] = r["expire"]
        else:
            merged[r["sku"]] = dict(r)
            order.append(r["sku"])

    rows = []
    for i, sku in enumerate(order, start=1):
        item = merged[sku]
        rows.append({
            "no": i,
            "sku": item["sku"],
            "name": item["name"],
            "box": LS.guess_box(settings, item["sku"], item["name"], item["qty"]),
            "qty": item["qty"],
            "expire": item["expire"],
        })

    return {
        "supplier": supplier or "",
        "center": _norm_center(center) if center else "",
        "center_raw": (center or "").strip(),
        "date": date or "",
        "pallet_total": pallet_total,
        "rows": rows,
    }


def parse_거래명세서(data: bytes, settings: dict | None = None) -> dict:
    """다중 PO PDF. (쿠팡제출용/공급사보관용)으로 페이지 중복 → PO 기준 dedup.

    2026-06부터 국문 양식('거래 명세서 / 쿠팡 제출용')으로 바뀌어
    먼저 국문 양식을 시도하고, 아니면 기존 영문 양식으로 파싱한다.
    SKU 별로 확정수량과 유통기한을 모으고, 확정수량 0 라인은 제외한다.
    """
    settings = settings or LS.get_settings("coupang")
    # 국문 신양식 감지
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        head = (pdf.pages[0].extract_text() or "") if pdf.pages else ""
    if "쿠팡 제출용" in head or ("발주번호" in head and "납품 센터" in head):
        return _parse_거래명세서_국문(data, settings)
    supplier = fc = date = None
    seen_po = set()
    ordered_rows: list[dict] = []

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pg in pdf.pages:
            txt = pg.extract_text() or ""
            if supplier is None:
                # '1. Supplier Info.' 헤더 줄이 아닌, 실제 'Supplier [코드] 상호' 줄
                m = re.search(r"^Supplier\s+(?:\[[^\]]+\]\s*)?(.+)$", txt, re.MULTILINE)
                if m:
                    supplier = m.group(1).strip()
            if fc is None:
                m = re.search(r"FC Name\s+([^\n]+)", txt)
                if m:
                    fc = m.group(1).strip()
            if date is None:
                m = re.search(r"Expected Receiving Date\s+([\d\-]+)", txt)
                if m:
                    date = m.group(1)

            mpo = re.search(r"PO no\.\s*(\d+)", txt)
            if not mpo:
                continue
            po = mpo.group(1)
            if po in seen_po:
                continue  # 중복 페이지(보관용)
            seen_po.add(po)

            for tbl in pg.extract_tables():
                for row in tbl:
                    if not row or len(row) < 4:
                        continue
                    c0 = (row[0] or "").strip()
                    c1 = (row[1] or "").strip()
                    # 데이터행: c0=번호(숫자), c1='SKUID\n상품명\n바코드...'
                    if not (c0.isdigit() and "\n" in c1):
                        continue
                    parts = c1.split("\n")
                    sku_id = parts[0].strip()
                    if not sku_id.isdigit():
                        continue
                    # 상품명은 여러 줄로 줄바꿈될 수 있음 → 바코드/옵션 줄 전까지 이어붙임
                    name_lines = []
                    for p in parts[1:]:
                        p = p.strip()
                        if not p:
                            continue
                        if p == "-" or re.fullmatch(r"[A-Z]?\d{9,}", p):
                            break  # 바코드/BOX바코드 줄
                        name_lines.append(p)
                    name = " ".join(name_lines)
                    qty = _to_int(row[3] if len(row) > 3 else 0)  # 확정수량
                    if qty <= 0:
                        continue  # 미확정/취소 라인 제외
                    # 유통기한: 행 전체에서 마지막 날짜(=Use-By Date)
                    dates = re.findall(r"\d{4}-\d{2}-\d{2}", " ".join(str(c) for c in row if c))
                    expire = dates[-1] if dates else ""
                    ordered_rows.append({
                        "sku": sku_id,
                        "name": _norm_product(name),
                        "qty": qty,
                        "expire": expire,
                    })

    # 동일 SKU 가 여러 PO 에 걸쳐 있으면 수량 합산(첫 등장 순서 유지)
    merged: dict[str, dict] = {}
    order: list[str] = []
    for r in ordered_rows:
        if r["sku"] in merged:
            merged[r["sku"]]["qty"] += r["qty"]
            if not merged[r["sku"]]["expire"]:
                merged[r["sku"]]["expire"] = r["expire"]
        else:
            merged[r["sku"]] = dict(r)
            order.append(r["sku"])

    rows = []
    for i, sku in enumerate(order, start=1):
        item = merged[sku]
        rows.append({
            "no": i,
            "sku": item["sku"],
            "name": item["name"],
            "box": LS.guess_box(settings, item["sku"], item["name"], item["qty"]),  # 입수 추정(수정 가능)
            "qty": item["qty"],
            "expire": item["expire"],
        })

    return {
        "supplier": supplier or "",
        "center": _norm_center(fc) if fc else "",
        "center_raw": (fc or "").strip(),
        "date": date or "",
        "rows": rows,
    }


# ── 부착리스트 파서 ────────────────────────────────────────────────────────
def parse_부착리스트(data: bytes) -> dict:
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        txt = (pdf.pages[0].extract_text() or "") if pdf.pages else ""

    # 신양식 (물류부착문서): "* 부착문서 - 팔레트 겉면에 부착해주세요 / 팔레트 4-1 / 밀크런ID ..."
    if "부착문서" in txt or "밀크런ID" in txt:
        mm = re.search(r"(\d{7,})\s+(\d{4}-\d{2}-\d{2})", txt)  # '10599773 2026-07-17'
        milkrun = mm.group(1) if mm else ""
        date = mm.group(2) if mm else ""
        mp = re.search(r"팔레트\s*(\d+)\s*-\s*(\d+)", txt)
        pallet_total = max(1, int(mp.group(1))) if mp else 1
        mc = re.search(r"([가-힣A-Za-z0-9]+)\(\d+\)", txt)  # '전라광주2(44)'
        center = _norm_center(mc.group(0)) if mc else ""
        return {
            "milkrun": milkrun,
            "center": center,
            "date": date,
            "pallet": f"{pallet_total}-1",
            "pallet_total": pallet_total,
        }

    milkrun = re.search(r"밀크런\s*번호\s*[:：]\s*(\d+)", txt)
    center = re.search(r"받는\s*사람\s*[:：]\s*([^\n]+)", txt)
    date = re.search(r"입고예정일자\s*[:：]\s*([\d\-]+)", txt)
    pallet = re.search(r"팔레트\s*수량\s*[:：]\s*([\d\s\-]+)", txt)
    raw = center.group(1).strip() if center else None
    # '팔레트 수량 : 5 - 1' → 첫 숫자가 총 팔레트 수(=슬라이드 장수)
    pallet_str = re.sub(r"\s+", "", pallet.group(1)) if pallet else ""
    pallet_total = 1
    if pallet:
        nums = re.findall(r"\d+", pallet.group(1))
        if nums:
            pallet_total = max(1, int(nums[0]))
    return {
        "milkrun": milkrun.group(1) if milkrun else "",
        "center": _norm_center(raw) if raw else "",
        "date": date.group(1) if date else "",
        "pallet": pallet_str,
        "pallet_total": pallet_total,
    }


def _sniff_kind(filename: str, data: bytes) -> str:
    """파일명 + 내용으로 거래명세서/부착리스트 구분."""
    fn = (filename or "")
    if "부착" in fn:
        return "부착리스트"
    if "명세" in fn or "동봉" in fn:   # 신양식 파일명: 물류동봉문서_*.pdf
        return "거래명세서"
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            head = (pdf.pages[0].extract_text() or "") if pdf.pages else ""
    except Exception:
        return "기타"
    # 신양식 마커를 먼저 확인 — 국문 거래명세서에도 '밀크런'(운송타입)이 있어
    # 구양식 규칙('밀크런'→부착)이 먼저 돌면 오분류된다.
    if "부착문서" in head:
        return "부착리스트"
    if "쿠팡 제출용" in head or ("발주번호" in head and "납품 센터" in head):
        return "거래명세서"
    if "팔레트 부착" in head or "밀크런" in head:
        return "부착리스트"
    if "Transaction" in head or "PO no" in head or "Statement" in head:
        return "거래명세서"
    return "기타"


# ── 파싱 엔드포인트 ────────────────────────────────────────────────────────
@router.post("/parse")
async def parse_docs(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(400, "거래명세서/부착리스트 PDF를 업로드해주세요.")

    명세: dict | None = None
    부착: dict | None = None
    errors: list[str] = []
    settings = LS.get_settings("coupang")

    for up in files:
        if not up.filename:
            continue
        data = await up.read()
        kind = _sniff_kind(up.filename, data)
        try:
            if kind == "부착리스트":
                부착 = parse_부착리스트(data)
            elif kind == "거래명세서":
                d = parse_거래명세서(data, settings)
                # 여러 거래명세서 파일이면 SKU 행 누적
                if 명세 is None:
                    명세 = d
                else:
                    base_no = len(명세["rows"])
                    for r in d["rows"]:
                        r = dict(r)
                        r["no"] = base_no + r["no"]
                        명세["rows"].append(r)
                    명세["supplier"] = 명세["supplier"] or d["supplier"]
                    명세["center"] = 명세["center"] or d["center"]
                    명세["date"] = 명세["date"] or d["date"]
            else:
                errors.append(f"{up.filename}: 종류를 인식하지 못했습니다.")
        except Exception as e:  # noqa: BLE001
            errors.append(f"{up.filename}: {e}")

    if 명세 is None:
        raise HTTPException(
            400,
            "거래명세서를 인식하지 못했습니다. 쿠팡 거래명세서 PDF 를 포함해주세요. "
            + ("(" + " / ".join(errors) + ")" if errors else ""),
        )

    date = 명세.get("date") or (부착.get("date") if 부착 else "")

    # 거래명세서 ↔ 부착리스트 정보 교차 확인
    warnings: list[str] = []
    if 부착:
        if 부착.get("center") and 명세.get("center") and 부착["center"] != 명세["center"]:
            warnings.append(
                f"납품센터 불일치 — 거래명세서 '{명세['center']}' ≠ 부착리스트 '{부착['center']}'"
            )
        if 부착.get("date") and 명세.get("date") and 부착["date"] != 명세["date"]:
            warnings.append(
                f"입고예정일 불일치 — 거래명세서 '{명세['date']}' ≠ 부착리스트 '{부착['date']}'"
            )

    # 팔레트 수: 부착리스트 → 거래명세서(신양식 '팔레트수량') → 1
    pallet_total = ((부착.get("pallet_total") if 부착 else 0)
                    or 명세.get("pallet_total", 0) or 1)
    return {
        "supplier": 명세.get("supplier", ""),
        "center": (부착.get("center") if 부착 else "") or 명세.get("center", ""),
        "date": date,
        "milkrun": (부착.get("milkrun") if 부착 else ""),
        "pallet": (부착.get("pallet") if 부착 else "") or f"{pallet_total}-1",
        "pallet_total": pallet_total,
        "rows": 명세.get("rows", []),
        "parse_errors": errors,
        "warnings": warnings,
        "has_부착": 부착 is not None,
        "settings": settings,
    }


# ── 적재 설정 (파레트당 박스 수·입수 규칙) ─────────────────────────────────
class LoadRule(BaseModel):
    match: str = ""       # 상품번호(전체 일치) 또는 상품명 키워드(포함)
    per_box: int = 1      # 1박스당 수량


class CoupangLoadSettings(BaseModel):
    pallet_cap: int = 112
    pallet_total_mode: str = "auto"   # auto | attach | fixed
    pallet_total_fixed: int = 1
    default_per_box: int = 9
    bundle_is_box: bool = True
    rules: list[LoadRule] = []


@router.get("/settings")
def get_load_settings():
    return LS.get_settings("coupang")


@router.post("/settings")
def save_load_settings(body: CoupangLoadSettings):
    return LS.save_settings("coupang", body.model_dump())


# ── PPT 생성 ───────────────────────────────────────────────────────────────
class LoadRow(BaseModel):
    no: int | None = None
    sku: str = ""
    name: str = ""
    box: int = 0
    qty: int = 0
    expire: str = ""
    pallet: int = 1          # 이 품목이 실릴 팔레트(=슬라이드) 번호


class GeneratePayload(BaseModel):
    supplier: str = ""
    center: str = ""
    date: str = ""          # 'YYYY-MM-DD' 또는 빈 값
    milkrun: str = ""
    pallet: str = "1-1"     # (호환용) 단일 팔레트일 때의 표기
    pallet_total: int = 1   # 총 팔레트 수 = 생성할 슬라이드 장수
    rows: list[LoadRow] = []


def _set_para_text(para, text: str):
    """첫 run 의 서식을 유지한 채 단락 텍스트를 통째로 교체."""
    runs = para.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        para.add_run().text = text


def _fill_cell(cell, text):
    cell.text = "" if text is None else str(text)
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.bold = True


def _shape_text(shape) -> str:
    return shape.text_frame.text if shape.has_text_frame else ""


def _duplicate_slide(prs, src):
    """src 슬라이드의 모든 도형을 복제한 새 슬라이드를 같은 레이아웃으로 추가."""
    new_slide = prs.slides.add_slide(src.slide_layout)
    # 레이아웃이 자동 추가한 플레이스홀더 제거
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    # 원본 도형(직사각형/텍스트박스/표) 그대로 복사
    for shp in src.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def _fill_slide(slide, payload: "GeneratePayload", pallet_label: str,
                box_total: int, rows: list, month: str, day: str):
    """한 슬라이드(팔레트)에 머리글 + 표를 채운다."""
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table and table_shape is None:
            table_shape = shape
            continue
        if not shape.has_text_frame:
            continue
        txt = _shape_text(shape)

        # 팔레트 번호 / 박스수량 줄
        if "박스수량" in txt:
            _set_para_text(
                shape.text_frame.paragraphs[0],
                f" {pallet_label}  /     박스수량.  (  {box_total} BOX)",
            )
            continue

        # 입고예정일/센터명, 업체명, 입고요청서번호 (3줄)
        if "입고예정일자" in txt:
            paras = shape.text_frame.paragraphs
            if len(paras) >= 1:
                _set_para_text(
                    paras[0],
                    f"입고예정일자. (   {month}  월    {day}  일  )   /    납품센터명 ({payload.center})",
                )
            if len(paras) >= 2:
                _set_para_text(paras[1], f"업체명.         (    {payload.supplier}        )")
            if len(paras) >= 3:
                _set_para_text(paras[2], f"입고요청서번호.      ({payload.milkrun})")
            continue

    if table_shape is None:
        raise HTTPException(500, "양식에서 표를 찾을 수 없습니다.")
    tbl = table_shape.table
    tbl_elem = tbl._tbl
    trs = tbl_elem.findall(qn("a:tr"))
    if len(trs) < 2:
        raise HTTPException(500, "양식 표 구조가 올바르지 않습니다.")

    template_tr = copy.deepcopy(trs[1])  # 첫 데이터행을 행 양식으로 사용
    for tr in trs[1:]:                   # 기존 데이터행 모두 제거(헤더만 남김)
        tbl_elem.remove(tr)
    for _ in rows:
        tbl_elem.append(copy.deepcopy(template_tr))

    for i, r in enumerate(rows, start=1):
        cells = tbl.rows[i].cells
        values = [i, r.sku, r.name, r.box, r.qty, r.expire]
        for ci, val in enumerate(values):
            if ci < len(cells):
                _fill_cell(cells[ci], val)


def _strip_revision_metadata(buf: io.BytesIO) -> io.BytesIO:
    """생성된 PPTX에서 공동편집/리비전 메타데이터를 제거한다.

    양식이 공동편집(co-authoring) 세션 중 저장돼 changesInfo/revisionInfo 를
    물고 있으면, 코드가 슬라이드를 복제·수정했을 때 이 변경이력과 실제 구성이
    어긋나 PowerPoint 가 병합/복구·편집제한 상태로 파일을 연다.
    이 파트들은 presentation.xml 본문이 아니라 rels 에만 매달린 고아 참조라
    제거해도 문서 편집에는 아무 영향이 없다.
    """
    import zipfile

    src = zipfile.ZipFile(buf)
    drop = {"ppt/revisionInfo.xml"}
    remove_names = {
        n for n in src.namelist()
        if n in drop or n.startswith("ppt/changesInfos/")
    }
    if not remove_names:
        buf.seek(0)
        return buf

    _rev_re = re.compile(r"changesInfo|revisionInfo")
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            if item.filename in remove_names:
                continue
            data = src.read(item.filename)
            if item.filename == "ppt/_rels/presentation.xml.rels":
                txt = data.decode("utf-8")
                # 관련 Relationship 항목 제거
                txt = re.sub(
                    r'<Relationship\b[^>]*Target="[^"]*(?:changesInfos/|revisionInfo)[^"]*"[^>]*/>',
                    "", txt,
                )
                data = txt.encode("utf-8")
            elif item.filename == "[Content_Types].xml":
                txt = data.decode("utf-8")
                txt = re.sub(
                    r'<Override\b[^>]*PartName="[^"]*(?:changesInfos/|revisionInfo)[^"]*"[^>]*/>',
                    "", txt,
                )
                data = txt.encode("utf-8")
            dst.writestr(item, data)
    out.seek(0)
    return out


def build_pptx(payload: GeneratePayload) -> io.BytesIO:
    if not os.path.exists(TEMPLATE_PATH):
        raise HTTPException(500, "적재리스트 양식 파일을 찾을 수 없습니다.")

    prs = Presentation(TEMPLATE_PATH)

    # 입고예정일 → 월/일
    month = day = ""
    m = re.search(r"(\d{4})-(\d{2})-(\d{2})", payload.date or "")
    if m:
        month, day = str(int(m.group(2))), str(int(m.group(3)))

    rows = payload.rows or []
    # 총 팔레트 수 = 지정값과 실제 사용된 최대 팔레트 번호 중 큰 값(최소 1)
    max_used = max((int(r.pallet or 1) for r in rows), default=1)
    total = max(int(payload.pallet_total or 1), max_used, 1)

    # 팔레트 번호별 그룹핑(범위를 벗어난 값은 1번으로)
    groups: dict[int, list] = {p: [] for p in range(1, total + 1)}
    for r in rows:
        p = int(r.pallet or 1)
        if p < 1 or p > total:
            p = 1
        groups[p].append(r)

    base = prs.slides[0]
    # 필요한 만큼 슬라이드 확보(원본 1장 + 복제). 복제는 채우기 전에 모두 수행.
    slides = [base] + [_duplicate_slide(prs, base) for _ in range(total - 1)]

    for idx, slide in enumerate(slides, start=1):
        grp = groups.get(idx, [])
        box_total = sum(_to_int(r.box) for r in grp)
        label = f"{total}-{idx}"
        _fill_slide(slide, payload, label, box_total, grp, month, day)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return _strip_revision_metadata(buf)


@router.post("/generate")
def generate_pptx(payload: GeneratePayload):
    buf = build_pptx(payload)
    center = _norm_center(payload.center) or "적재리스트"
    fname = f"{center}_적재리스트.pptx"
    ascii_fallback = "load_list.pptx"
    from urllib.parse import quote
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": (
                f"attachment; filename={ascii_fallback}; "
                f"filename*=UTF-8''{quote(fname)}"
            )
        },
    )
