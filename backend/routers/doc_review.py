# -*- coding: utf-8 -*-
"""쿠팡 밀크런 서류 자동분류 + 규칙기반 교차검토 (API 0원)

여러 서류(거래명세서/부착리스트/적재리스트/출고요청 엑셀)를 센터 구분 없이
한꺼번에 업로드하면, 내용에서 센터·밀크런번호를 추출해 센터별로 묶고,
거래명세서(쿠팡 확정수량) ↔ 적재리스트 ↔ 출고요청 수량/박스를 코드로 대조한다.
외부 API를 호출하지 않는다.
"""
from fastapi import APIRouter, UploadFile, File, HTTPException
from html.parser import HTMLParser
import html
import io
import re
import zipfile
from collections import defaultdict

import pdfplumber
from pptx import Presentation
import pandas as pd

from routers import kurly_label, coupang_load

router = APIRouter()


# ── 정규화 유틸 ────────────────────────────────────────────────────────────
def norm_center(name: str) -> str:
    """'창원1(18)' / '대구7(RC)' / '창원1' → '창원1'"""
    if not name:
        return ""
    s = str(name).strip()
    s = re.sub(r"\s*\(.*?\)\s*$", "", s)   # 끝 괄호 제거
    s = re.sub(r"\s+", "", s)
    return s


def norm_product(name: str) -> str:
    if not name:
        return ""
    return re.sub(r"\s+", " ", str(name)).strip()


def pkey(name: str) -> str:
    """상품명 조인키: 모든 공백 제거 (예: '3개 입' == '3개입')"""
    return re.sub(r"\s+", "", str(name or "")).lower()


def _to_int(v) -> int:
    try:
        return int(re.sub(r"[^\d\-]", "", str(v)) or 0)
    except Exception:
        return 0


# ── 서류 종류 판별 ─────────────────────────────────────────────────────────
def classify(fname: str) -> str:
    n = (fname or "").lower()
    if n.endswith((".xlsx", ".xls")):
        return "출고요청"
    if "거래명세서" in fname or "명세서" in fname or "동봉" in fname:  # 신양식: 물류동봉문서_*.pdf
        return "거래명세서"
    if "부착리스트" in fname or "부착" in fname:
        return "부착리스트"
    if "적재리스트" in fname or "적재" in fname:
        return "적재리스트"
    if "번들" in fname:
        return "번들라벨"
    return "기타"


# ── HTML 표(.xls로 위장한 HTML) 파서 ────────────────────────────────────────
class _HTMLTable(HTMLParser):
    def __init__(self):
        super().__init__()
        self.rows: list[list[str]] = []
        self._cur: list[str] | None = None
        self._buf: list[str] | None = None

    def handle_starttag(self, tag, attrs):
        if tag == "tr":
            self._cur = []
        elif tag in ("td", "th"):
            self._buf = []

    def handle_endtag(self, tag):
        if tag == "tr" and self._cur is not None:
            self.rows.append(self._cur)
            self._cur = None
        elif tag in ("td", "th") and self._buf is not None and self._cur is not None:
            self._cur.append(re.sub(r"\s+", " ", "".join(self._buf)).strip())
            self._buf = None

    def handle_data(self, d):
        if self._buf is not None:
            self._buf.append(d)


def _is_html_xls(data: bytes) -> bool:
    head = data[:512].lstrip().lower()
    return head.startswith(b"<") or b"<html" in head or b"<table" in head


def _excel_rows(data: bytes) -> list:
    """진짜 엑셀(.xlsx 등)을 문자열 2차원 배열로 — 접수내역을 엑셀에서
    다시 저장/편집한 경우 지원. 숫자는 정수 문자열로, 날짜는 그대로 str."""
    try:
        raw = pd.read_excel(io.BytesIO(data), header=None)
    except Exception:
        return []

    def _cell(v):
        if v is None or (isinstance(v, float) and pd.isna(v)):
            return ""
        if isinstance(v, float) and v.is_integer():
            return str(int(v))  # 10612861.0 → '10612861'
        return str(v).strip()

    return [[_cell(v) for v in raw.iloc[i].tolist()] for i in range(len(raw))]


def parse_밀크런접수(data: bytes) -> dict:
    """쿠팡 밀크런 접수내역 목록 — HTML로 위장한 원본 .xls 또는
    엑셀에서 다시 저장한 진짜 .xlsx 모두 지원.

    밀크런번호별 도착센터·박스수·팔레트수·발주번호(PO)를 담은 마스터 목록.
    반환: { 밀크런번호: {milkrun, status, center, center_raw, boxes, pallets, pos:[...]} }
    """
    if _is_html_xls(data):
        p = _HTMLTable()
        p.feed(data.decode("utf-8", "ignore"))
        rows = [r for r in p.rows if r]
    else:
        rows = _excel_rows(data)
    if not rows:
        return {}
    # 헤더 행 탐색 — 엑셀 편집으로 위에 빈 행/제목이 생겨도 찾도록 상위 몇 줄 스캔
    hdr_i = next((i for i, r in enumerate(rows[:10])
                  if any("밀크런번호" in str(c) for c in r)), None)
    if hdr_i is None:
        return {}
    header = rows[hdr_i]
    rows = rows[hdr_i:]
    idx = {h: i for i, h in enumerate(header)}

    def col(row, name):
        i = idx.get(name)
        return row[i].strip() if (i is not None and i < len(row)) else ""

    out: dict[str, dict] = {}
    for row in rows[1:]:
        mk = col(row, "밀크런번호")
        if not mk or not mk.isdigit():
            continue
        raw_center = col(row, "물류센터")
        pos = [p for p in re.split(r"[\s/]+", col(row, "발주번호")) if p.isdigit()]
        out[mk] = {
            "milkrun": mk,
            "status": col(row, "상태") or "정상",
            "center_raw": raw_center,
            "center": norm_center(raw_center),
            "boxes": _to_int(col(row, "박스수")),
            "pallets": _to_int(col(row, "총 팔레트 수량")),
            "date": col(row, "픽업일")[:10],  # 엑셀 재저장 시 '2026-07-14 00:00:00' 방지
            "pos": pos,
        }
    return out


def is_밀크런접수(data: bytes) -> bool:
    """엑셀 바이트가 밀크런 접수내역 목록인지(밀크런번호·물류센터 컬럼 보유) 판별."""
    if _is_html_xls(data):
        head = data[:4096].decode("utf-8", "ignore")
        return "밀크런번호" in head and ("물류센터" in head or "발주번호" in head)
    if data[:2] == b"PK":  # 진짜 엑셀(.xlsx) — 상위 행에서 헤더 확인
        try:
            raw = pd.read_excel(io.BytesIO(data), header=None, nrows=10)
        except Exception:
            return False
        joined = "|".join(str(v) for v in raw.fillna("").values.flatten())
        return "밀크런번호" in joined and ("물류센터" in joined or "발주번호" in joined)
    return False


# ── 파서 (bytes 입력) ──────────────────────────────────────────────────────
def parse_부착리스트(data: bytes) -> dict:
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        txt = pdf.pages[0].extract_text() or "" if pdf.pages else ""

    # 신양식 (물류부착문서): "* 부착문서 - 팔레트 겉면에 부착해주세요 / 팔레트 4-1 / 밀크런ID ..."
    if "부착문서" in txt or "밀크런ID" in txt:
        mm = re.search(r"(\d{7,})\s+(\d{4}-\d{2}-\d{2})", txt)  # '10599773 2026-07-17'
        mp = re.search(r"팔레트\s*(\d+)\s*-\s*(\d+)", txt)
        mc = re.search(r"([가-힣A-Za-z0-9]+)\(\d+\)", txt)      # '창원4(06)'
        raw = mc.group(0) if mc else None
        return {
            "milkrun": mm.group(1) if mm else None,
            "center": norm_center(raw) if raw else None,
            "center_raw": raw,
            "date": mm.group(2) if mm else None,
            "pallet": f"{mp.group(1)}-{mp.group(2)}" if mp else None,
        }

    milkrun = re.search(r"밀크런\s*번호\s*[:：]\s*(\d+)", txt)
    center = re.search(r"받는\s*사람\s*[:：]\s*([^\n]+)", txt)
    date = re.search(r"입고예정일자\s*[:：]\s*([\d\-]+)", txt)
    pallet = re.search(r"팔레트\s*수량\s*[:：]\s*([\d\s\-]+)", txt)
    raw = center.group(1).strip() if center else None
    return {
        "milkrun": milkrun.group(1) if milkrun else None,
        "center": norm_center(raw) if raw else None,
        "center_raw": raw,
        "date": date.group(1) if date else None,
        "pallet": pallet.group(1).strip() if pallet else None,
    }


def parse_적재리스트(data: bytes) -> dict:
    prs = Presentation(io.BytesIO(data))
    center = milkrun = date = None
    boxtotal = None
    rows = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                m = re.search(r"납품센터명\s*\(([^)]+)\)", t)
                if m:
                    center = norm_center(m.group(1))
                m = re.search(r"입고요청서번호\.?\s*\(([0-9]+)\)", t)
                if m:
                    milkrun = m.group(1)
                m = re.search(r"입고예정일자\.?\s*\(\s*(\d+)\s*월\s*(\d+)\s*일", t)
                if m:
                    date = f"{int(m.group(1)):02d}-{int(m.group(2)):02d}"
                m = re.search(r"박스수량\.?\s*\(\s*([\d,]+)\s*BOX", t)
                if m:
                    boxtotal = (boxtotal or 0) + int(m.group(1).replace(",", ""))
            if shape.has_table:
                tbl = shape.table
                for r in list(tbl.rows)[1:]:
                    cells = [c.text.strip() for c in r.cells]
                    # NO, 상품번호, 상품명, BOX수량, 수량, 유통기한
                    if len(cells) < 5 or not cells[0]:
                        continue
                    rows.append({
                        "sku": cells[1].strip(),
                        "name": norm_product(cells[2]),
                        "box": _to_int(cells[3]),
                        "qty": _to_int(cells[4]),
                    })
    return {"center": center, "milkrun": milkrun, "date": date,
            "boxtotal": boxtotal, "rows": rows}


def _parse_거래명세서_국문(data: bytes) -> dict:
    """2026-06 이후 쿠팡 국문 양식 — '거래 명세서 / 쿠팡 제출용' (물류동봉문서).

    페이지가 (쿠팡 제출용/업체 보관용) 쌍으로 반복 → 발주번호 기준 dedup.
    발주정보 표에서 납품 센터·도착예정일, 상품 표에서 상품번호·상품명·확정수량 추출.
    """
    fc = None
    date = None
    pos = {}    # po -> {sku: qty(확정)}
    names = {}  # sku -> name
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pg in pdf.pages:
            txt = pg.extract_text() or ""
            mpo = re.search(r"발주번호\s*(\d+)", txt)
            if not mpo:
                continue
            po = mpo.group(1)
            if po in pos:
                continue  # 업체 보관용(중복) 페이지

            kv: dict = {}
            item_tbl = None
            for tbl in pg.extract_tables() or []:
                if not tbl:
                    continue
                header = [(c or "").strip() for c in tbl[0]]
                if (any(("상품번호" in h or "상품코드" in h) for h in header)
                        and any("발주수량" in h for h in header)):
                    item_tbl = tbl
                    continue
                for row in tbl:
                    if row and len(row) >= 2 and row[0]:
                        kv[str(row[0]).strip()] = str(row[1] or "").strip()

            if not fc:
                c = kv.get("납품 센터") or kv.get("납품센터") or kv.get("FC명")
                if c:
                    fc = norm_center(c)
            if not date:
                raw = re.sub(r"\D", "", kv.get("물류센터 도착예정일") or kv.get("입고예정일") or "")
                if len(raw) >= 8:
                    date = f"{raw[:4]}-{raw[4:6]}-{raw[6:8]}"

            skus: dict = {}
            if item_tbl:
                header = [(c or "").replace("\n", "") for c in item_tbl[0]]
                i_no = next((i for i, h in enumerate(header) if h.strip() == "No"), 0)
                i_sku = next((i for i, h in enumerate(header)
                              if "상품번호" in h or "상품코드" in h), 1)
                i_name = next((i for i, h in enumerate(header) if "상품명" in h), i_sku)
                i_qty = next((i for i, h in enumerate(header) if "발주수량" in h), 3)
                i_recv = next((i for i, h in enumerate(header) if "확정수량" in h), None)
                for row in item_tbl[1:]:
                    if not row or len(row) <= max(i_sku, i_qty):
                        continue
                    c0 = str(row[i_no] or "").strip()
                    if not c0.isdigit():
                        continue  # '합계' 등
                    cell_lines = [p.strip() for p in str(row[i_sku] or "").split("\n") if p.strip()]
                    if i_sku == i_name:
                        # 결합 셀: '상품코드\n상품명...\nBarcode\nBOX바코드'
                        if not cell_lines or not cell_lines[0].isdigit():
                            continue
                        sku_id = cell_lines[0]
                        name_src = cell_lines[1:]
                    else:
                        sku_id = cell_lines[0] if cell_lines else ""
                        name_src = [p.strip() for p in str(row[i_name] or "").split("\n") if p.strip()]
                    if not sku_id.isdigit():
                        continue
                    nm_lines = []
                    for p in name_src:
                        if p == "-" or re.fullmatch(r"[A-Z]?\d{9,}", p):
                            break
                        nm_lines.append(p)
                    # 확정수량 (없는 변형이면 발주수량)
                    if i_recv is not None:
                        qty = _to_int(row[i_recv] if i_recv < len(row) else 0)
                    else:
                        qty = _to_int(row[i_qty])
                    skus[sku_id] = skus.get(sku_id, 0) + qty
                    names[sku_id] = norm_product(" ".join(nm_lines))
            pos[po] = skus
    return {"fc": fc, "date": date, "pos": pos, "names": names}


def parse_거래명세서(data: bytes) -> dict:
    """다중 PO PDF. (쿠팡제출용/공급사보관용)으로 페이지 중복 → PO 기준 dedup.

    2026-06부터 국문 양식('거래 명세서 / 쿠팡 제출용')으로 바뀌어
    먼저 국문 양식을 시도하고, 아니면 기존 영문 양식으로 파싱한다.
    """
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        head = (pdf.pages[0].extract_text() or "") if pdf.pages else ""
    if "쿠팡 제출용" in head or ("발주번호" in head and ("납품 센터" in head or "FC명" in head)):
        return _parse_거래명세서_국문(data)

    fc = None
    date = None
    pos = {}   # po -> {sku: qty}
    names = {}  # sku -> name
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pg in pdf.pages:
            txt = pg.extract_text() or ""
            mpo = re.search(r"PO no\.\s*(\d+)", txt)
            mfc = re.search(r"FC Name\s+([^\n]+)", txt)
            mdate = re.search(r"Expected Receiving Date\s+([\d\-]+)", txt)
            if mfc and not fc:
                fc = norm_center(mfc.group(1))
            if mdate and not date:
                date = mdate.group(1)
            if not mpo:
                continue
            po = mpo.group(1)
            if po in pos:
                continue   # 중복 페이지
            skus = {}
            for tbl in pg.extract_tables():
                for row in tbl:
                    if not row or len(row) < 4:
                        continue
                    c0 = (row[0] or "").strip()
                    c1 = (row[1] or "").strip()
                    # 데이터행: c0=번호(숫자), c1='SKUID\n상품명\n바코드...'
                    if c0.isdigit() and "\n" in c1:
                        parts = c1.split("\n")
                        sku_id = parts[0].strip()
                        nm = parts[1].strip() if len(parts) > 1 else ""
                        conf = (row[3] or row[2] or "").strip()
                        qty = _to_int(conf)
                        if sku_id.isdigit():
                            skus[sku_id] = skus.get(sku_id, 0) + qty
                            names[sku_id] = norm_product(nm)
            pos[po] = skus
    return {"fc": fc, "date": date, "pos": pos, "names": names}


def parse_출고요청(data: bytes) -> dict:
    """모든 센터를 한 장에 담은 마스터 엑셀. center -> [{name, qty, box}]"""
    xl = pd.ExcelFile(io.BytesIO(data))
    df = xl.parse(xl.sheet_names[0])
    out = defaultdict(list)
    for _, r in df.iterrows():
        center = norm_center(r.get("수취인"))
        if not center:
            continue
        qty = _to_int(r.get("수량"))
        # 박스수는 '주문메모' 컬럼에 들어있는 양식. 안되면 '택배사' 폴백.
        box = 0
        for col in ("주문메모", "택배사"):
            v = _to_int(r.get(col))
            if v:
                box = v
                break
        out[center].append({"name": norm_product(r.get("상품명")), "qty": qty, "box": box})
    return dict(out)


# ── 교차검토 ───────────────────────────────────────────────────────────────
ALL_DOCS = ["거래명세서", "부착리스트", "적재리스트", "출고요청"]


def _worst(*statuses: str) -> str:
    order = {"ok": 0, "warn": 1, "error": 2}
    return max(statuses, key=lambda s: order.get(s, 0)) if statuses else "ok"


def build_group(center: str, 명세: dict, 부착: dict, 적재: dict, req: list,
                mk_index: dict | None = None) -> dict:
    present = [k for k, v in [("거래명세서", 명세), ("부착리스트", 부착),
                              ("적재리스트", 적재), ("출고요청", req)] if v]
    missing = [k for k in ALL_DOCS if k not in present]

    checks = []
    statuses = []

    # 누락 서류
    if missing:
        checks.append({"label": "서류 구비", "status": "warn",
                       "detail": f"누락: {', '.join(missing)}"})
        statuses.append("warn")
    else:
        checks.append({"label": "서류 구비", "status": "ok", "detail": "4종 모두 있음"})

    # 밀크런번호 일치 (부착 ↔ 적재)
    mks = sorted({x for x in [부착 and 부착.get("milkrun"),
                              적재 and 적재.get("milkrun")] if x})
    if len(mks) == 1:
        checks.append({"label": "밀크런번호 일치", "status": "ok", "detail": mks[0]})
    elif len(mks) == 0:
        checks.append({"label": "밀크런번호 일치", "status": "warn", "detail": "정보 없음"})
        statuses.append("warn")
    else:
        checks.append({"label": "밀크런번호 일치", "status": "error",
                       "detail": " ≠ ".join(mks)})
        statuses.append("error")

    # 센터명 일치 (3종)
    cset = sorted({x for x in [부착 and 부착.get("center"), 적재 and 적재.get("center"),
                               명세 and 명세.get("fc")] if x})
    if len(cset) <= 1:
        checks.append({"label": "센터명 일치", "status": "ok", "detail": center})
    else:
        checks.append({"label": "센터명 일치", "status": "error", "detail": " ≠ ".join(cset)})
        statuses.append("error")

    # 입고예정일 일치 (부착 ↔ 명세)
    dset = sorted({x for x in [부착 and 부착.get("date"), 명세 and 명세.get("date")] if x})
    if len(dset) == 1:
        checks.append({"label": "입고예정일 일치", "status": "ok", "detail": dset[0]})
    elif len(dset) == 0:
        checks.append({"label": "입고예정일 일치", "status": "warn", "detail": "정보 없음"})
        statuses.append("warn")
    else:
        checks.append({"label": "입고예정일 일치", "status": "error", "detail": " ≠ ".join(dset)})
        statuses.append("error")

    # ── SKU 통합 대조 ──
    # 적재리스트 기준(상품번호+상품명 보유): sku -> [qty, box, name]
    agg적재 = defaultdict(lambda: [0, 0, ""])
    for r in (적재["rows"] if 적재 else []):
        agg적재[r["sku"]][0] += r["qty"]
        agg적재[r["sku"]][1] += r["box"]
        agg적재[r["sku"]][2] = r["name"]
    # 거래명세서 합계 by sku
    agg명세 = defaultdict(int)
    명세_names = 명세.get("names", {}) if 명세 else {}
    if 명세:
        for po, skus in 명세["pos"].items():
            for sku, q in skus.items():
                agg명세[sku] += q
    # 출고요청 by pkey(공백무시 상품명)
    agg요청 = defaultdict(lambda: [0, 0])
    for r in (req or []):
        agg요청[pkey(r["name"])][0] += r["qty"]
        agg요청[pkey(r["name"])][1] += r["box"]

    items = []
    tot = {"req": 0, "load": 0, "stmt": 0, "box_req": 0, "box_load": 0}
    sku_error = False
    for sku in sorted(set(agg적재) | set(agg명세)):
        load_q, load_b, nm = agg적재.get(sku, [0, 0, ""])
        # 적재리스트에 없는(명세 전용) SKU는 거래명세서의 상품명을 사용
        if not nm:
            nm = 명세_names.get(sku, "")
        stmt_q = agg명세.get(sku, 0)
        req_q, req_b = agg요청.get(pkey(nm), [0, 0]) if nm else (0, 0)
        if load_q == 0 and stmt_q == 0 and req_q == 0:
            continue   # 취소라인
        # 핵심: 적재 == 명세(쿠팡 확정)  → 불일치 시 error
        # 요청 != 적재 (요청 정보 있을 때) → warn
        if load_q != stmt_q:
            st = "error"
            sku_error = True
        elif req_q and req_q != load_q:
            st = "warn"
            statuses.append("warn")
        else:
            st = "ok"
        items.append({
            "sku": sku, "name": nm or "(이름 미상)",
            "req": req_q, "load": load_q, "stmt": stmt_q,
            "box": load_b, "status": st,
        })
        tot["req"] += req_q
        tot["load"] += load_q
        tot["stmt"] += stmt_q
        tot["box_req"] += req_b
        tot["box_load"] += load_b
    if sku_error:
        statuses.append("error")

    # 총 수량 체크
    qset = {tot["load"], tot["stmt"]}
    if tot["req"]:
        qset.add(tot["req"])
    if len(qset) == 1:
        checks.append({"label": "총 수량 (요청/적재/명세)", "status": "ok",
                       "detail": f"{tot['req']} / {tot['load']} / {tot['stmt']}"})
    else:
        st = "error" if tot["load"] != tot["stmt"] else "warn"
        checks.append({"label": "총 수량 (요청/적재/명세)", "status": st,
                       "detail": f"{tot['req']} / {tot['load']} / {tot['stmt']}"})
        statuses.append(st)

    # 총 박스 체크
    box표기 = 적재.get("boxtotal") if 적재 else None
    box_parts = []
    if tot["box_req"]:
        box_parts.append(f"요청 {tot['box_req']}")
    box_parts.append(f"적재 {tot['box_load']}")
    if box표기 is not None:
        box_parts.append(f"적재표기 {box표기}")
    box_vals = {tot["box_load"]}
    if tot["box_req"]:
        box_vals.add(tot["box_req"])
    if box표기 is not None:
        box_vals.add(box표기)
    box_st = "ok" if len(box_vals) == 1 else "warn"
    checks.append({"label": "총 박스", "status": box_st, "detail": " / ".join(box_parts)})
    if box_st != "ok":
        statuses.append("warn")

    # ── 밀크런 접수내역 대조 (목록 업로드 시에만) ──
    mk_matched = None
    if mk_index is not None:
        mk_no = mks[0] if len(mks) == 1 else None
        entry = mk_index.get(mk_no) if mk_no else None
        if entry is None and mk_no is None:
            # 서류에서 밀크런번호를 못 읽었을 때만 센터명으로 폴백 (정상 접수건 중).
            # 번호를 아는데 목록에 없으면 폴백하지 않는다 — 다른 건을 '정상 접수'로 오표시 방지.
            cands = [e for e in mk_index.values()
                     if e["center"] == center and e.get("status") != "취소"]
            entry = cands[0] if len(cands) == 1 else None

        # 발주번호(PO) 폴백 — 재접수 등으로 밀크런번호가 바뀌어도 발주번호가 겹치면 같은 건
        matched_by_po = False
        if entry is None and 명세:
            doc_pos = set((명세.get("pos") or {}).keys())
            if doc_pos:
                hits = [e for e in mk_index.values()
                        if doc_pos & set(e.get("pos") or []) and e.get("status") != "취소"]
                if len(hits) == 1:
                    entry = hits[0]
                    matched_by_po = True

        if entry is None:
            # 왜 못 찾았는지 진단 — 접수내역의 픽업일과 서류의 입고예정일이 다르면
            # (지난 주차 접수내역을 올린 경우가 대부분) 그 사실을 그대로 알려준다.
            list_dates = sorted({e.get("date", "") for e in mk_index.values() if e.get("date")})
            doc_date = dset[0] if len(dset) == 1 else ""
            detail = f"접수 목록({len(mk_index)}건)에 밀크런 {mk_no or '?'} 없음"
            if doc_date and list_dates and doc_date not in list_dates:
                detail += (f" — 올린 접수내역은 픽업일 {', '.join(list_dates)}만 담고 있어요. "
                           f"입고예정일 {doc_date}이 포함된 접수내역을 추가로 올려주세요 (여러 파일 동시 업로드 가능)")
            checks.append({"label": "밀크런 접수내역", "status": "warn", "detail": detail})
            statuses.append("warn")
        else:
            mk_matched = entry["milkrun"]
            # PO로 찾았는데 서류의 밀크런번호와 접수번호가 다르면 명시적으로 경고
            if matched_by_po and mk_no and entry["milkrun"] != mk_no:
                checks.append({"label": "밀크런 접수내역", "status": "warn",
                               "detail": (f"서류 번호({mk_no})와 접수번호({entry['milkrun']})가 다르지만 "
                                          f"발주번호가 일치 — 재접수/번호 변경 여부 확인하세요")})
                statuses.append("warn")
            # 접수 상태 (정상/취소)
            if entry.get("status") == "취소":
                checks.append({"label": "밀크런 접수 상태", "status": "error",
                               "detail": f"취소된 접수건 (밀크런 {entry['milkrun']})"})
                statuses.append("error")
            else:
                checks.append({"label": "밀크런 접수 상태", "status": "ok",
                               "detail": f"정상 접수 (밀크런 {entry['milkrun']})"})
            # 접수 센터 일치
            if entry["center"] and center and entry["center"] != center:
                checks.append({"label": "접수 센터 일치", "status": "error",
                               "detail": f"{center} ≠ 접수 {entry['center_raw']}"})
                statuses.append("error")
            else:
                checks.append({"label": "접수 센터 일치", "status": "ok",
                               "detail": entry["center_raw"] or center})
            # 접수 박스수 일치 (적재 총박스 기준)
            recv_box = entry.get("boxes") or 0
            if recv_box:
                cmp_box = box표기 if box표기 is not None else tot["box_load"]
                bst = "ok" if recv_box == cmp_box else "error"
                checks.append({"label": "접수 박스수 일치", "status": bst,
                               "detail": f"접수 {recv_box} / 서류 {cmp_box}"})
                if bst != "ok":
                    statuses.append(bst)
            # 발주번호(PO) 일치 (거래명세서 PO ↔ 접수 발주번호)
            recv_pos = set(entry.get("pos") or [])
            stmt_pos = set(명세["pos"].keys()) if 명세 else set()
            if recv_pos and stmt_pos:
                if recv_pos == stmt_pos:
                    checks.append({"label": "발주번호(PO) 일치", "status": "ok",
                                   "detail": f"{len(recv_pos)}건 일치"})
                else:
                    only_recv = sorted(recv_pos - stmt_pos)
                    only_stmt = sorted(stmt_pos - recv_pos)
                    parts = []
                    if only_recv:
                        parts.append(f"접수만: {', '.join(only_recv)}")
                    if only_stmt:
                        parts.append(f"명세만: {', '.join(only_stmt)}")
                    checks.append({"label": "발주번호(PO) 일치", "status": "error",
                                   "detail": " / ".join(parts)})
                    statuses.append("error")

    status = _worst(*statuses) if statuses else "ok"
    return {
        "center": center,
        "milkrun": (mks[0] if len(mks) == 1 else mk_matched),
        "date": (dset[0] if len(dset) == 1 else None),
        "status": status,
        "present": present,
        "missing": missing,
        "checks": checks,
        "items": items,
        "totals": tot,
    }


# ── 엔드포인트 ─────────────────────────────────────────────────────────────
@router.post("/coupang")
async def review_coupang(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(400, "파일을 하나 이상 업로드해주세요.")

    출고요청: dict = {}
    bun = defaultdict(dict)   # center -> {부착, 적재, 명세}
    mk_index: dict = {}       # 밀크런번호 -> 접수내역
    has_milkrun_list = False
    errors = []
    parsed_files = []

    for up in files:
        if not up.filename:
            continue
        kind = classify(up.filename)
        data = await up.read()
        try:
            if kind == "출고요청":
                # 밀크런 접수내역 목록(HTML .xls)인지 먼저 판별
                if is_밀크런접수(data) or "milkrun" in up.filename.lower():
                    idx = parse_밀크런접수(data)
                    mk_index.update(idx)
                    has_milkrun_list = True
                    # 픽업일을 배지에 표시 — 지난 주차 접수내역을 올렸을 때 바로 보이게
                    _dates = sorted({e.get("date", "") for e in idx.values() if e.get("date")})
                    _tag = f"밀크런접수({len(idx)}건"
                    if _dates:
                        _tag += f", 픽업일 {'·'.join(_dates)}"
                    parsed_files.append({"name": up.filename, "kind": _tag + ")"})
                    continue
                # 택배 출고요청은 밀크런 검토 대상 아님
                if "택배" in up.filename:
                    parsed_files.append({"name": up.filename, "kind": "택배(제외)"})
                    continue
                출고요청.update(parse_출고요청(data))
            elif kind == "부착리스트":
                d = parse_부착리스트(data)
                if d["center"]:
                    bun[d["center"]]["부착"] = d
            elif kind == "적재리스트":
                d = parse_적재리스트(data)
                if d["center"]:
                    bun[d["center"]]["적재"] = d
            elif kind == "거래명세서":
                d = parse_거래명세서(data)
                if d["fc"]:
                    bun[d["fc"]]["명세"] = d
            else:
                parsed_files.append({"name": up.filename, "kind": kind})
                continue
            parsed_files.append({"name": up.filename, "kind": kind})
        except Exception as e:
            errors.append(f"{up.filename}: {e}")

    centers = sorted(set([c for c in bun if c]) | set(출고요청.keys()))
    if not centers:
        raise HTTPException(
            400,
            "인식 가능한 쿠팡 밀크런 서류를 찾지 못했습니다. "
            "(거래명세서/부착리스트/적재리스트 PDF·PPTX, 출고요청 엑셀)"
        )

    mk_arg = mk_index if has_milkrun_list else None
    groups = []
    for c in centers:
        b = bun.get(c, {})
        groups.append(build_group(
            c, b.get("명세"), b.get("부착"), b.get("적재"), 출고요청.get(c), mk_arg
        ))

    # 접수는 됐는데(정상) 서류가 하나도 없는 밀크런 건을 별도로 노출
    if has_milkrun_list:
        covered = {g["milkrun"] for g in groups if g.get("milkrun")}
        covered_centers = {g["center"] for g in groups}
        for mk, e in sorted(mk_index.items()):
            if e.get("status") == "취소":
                continue
            if mk in covered or e["center"] in covered_centers:
                continue
            groups.append({
                "center": e["center_raw"] or e["center"],
                "milkrun": mk,
                "date": e.get("date"),
                "status": "error",
                "present": [],
                "missing": ALL_DOCS,
                "checks": [{
                    "label": "밀크런 접수내역", "status": "error",
                    "detail": f"접수됨(박스 {e.get('boxes', 0)}, PO {', '.join(e.get('pos') or []) or '-'}) — 서류 없음",
                }],
                "items": [],
                "totals": {"req": 0, "load": 0, "stmt": 0, "box_req": 0, "box_load": 0},
            })

    summary = {
        "total": len(groups),
        "ok": sum(1 for g in groups if g["status"] == "ok"),
        "warn": sum(1 for g in groups if g["status"] == "warn"),
        "error": sum(1 for g in groups if g["status"] == "error"),
    }
    return {
        "groups": groups,
        "summary": summary,
        "files": parsed_files,
        "parse_errors": errors,
    }


# ── 메일용 [출고 수량 상세] 자동 생성 ────────────────────────────────────────
def _display_center(raw: str) -> str:
    """센터 표기 정리: 숫자 코드 괄호만 제거하고 (RC) 등 문자 표기는 유지.

    '인천4(38)' → '인천4',  '인천41(RC)' → '인천41(RC)',  'XRC11(RC)' → 그대로
    """
    s = (raw or "").strip()
    s = re.sub(r"\s*\(\s*\d+\s*\)\s*$", "", s)   # 끝의 (숫자) 만 제거
    return re.sub(r"\s+", " ", s).strip()


def _center_sortkey(name: str):
    """'인천4' < '인천14' < '인천32' 처럼 문자+숫자 자연 정렬."""
    m = re.match(r"^(.*?)(\d+)?$", norm_center(name))
    prefix = m.group(1) if m else name
    num = int(m.group(2)) if (m and m.group(2)) else -1
    return (prefix, num)


@router.post("/coupang-mail")
async def coupang_mail_breakdown(files: list[UploadFile] = File(...)):
    """거래명세서(들)를 업로드하면 센터별 [출고 수량 상세] 텍스트를 만들어 준다.

    확정수량(Vendor Confirmed Qty)·박스(입수 추정)·유통기한(Use-By Date)을
    거래명세서에서 뽑아 센터별로 묶는다. 메일 본문은 사용자가 직접 작성한다.
    """
    if not files:
        raise HTTPException(400, "거래명세서 PDF를 하나 이상 업로드해주세요.")

    # center_raw -> {"center": norm, "rows": [{sku,name,qty,box,expire}], "date":..}
    centers: dict[str, dict] = {}
    errors: list[str] = []
    parsed_files: list[dict] = []

    for up in files:
        if not up.filename:
            continue
        data = await up.read()
        if not up.filename.lower().endswith(".pdf"):
            parsed_files.append({"name": up.filename, "kind": "제외(비 PDF)"})
            continue
        kind = coupang_load._sniff_kind(up.filename, data)
        if kind != "거래명세서":
            parsed_files.append({"name": up.filename, "kind": f"제외({kind})"})
            continue
        try:
            d = coupang_load.parse_거래명세서(data)
        except Exception as e:  # noqa: BLE001
            errors.append(f"{up.filename}: {e}")
            continue
        disp = _display_center(d.get("center_raw") or d.get("center")) or "(센터 미상)"
        g = centers.setdefault(disp, {"center": d.get("center", ""),
                                      "date": d.get("date", ""), "rows": []})
        # 같은 센터 명세가 여러 장이면 SKU 합산
        by_sku = {r["sku"]: r for r in g["rows"]}
        for r in d["rows"]:
            if r["sku"] in by_sku:
                ex = by_sku[r["sku"]]
                ex["qty"] += r["qty"]
                ex["box"] += r["box"]
                if not ex.get("expire"):
                    ex["expire"] = r.get("expire", "")
            else:
                nr = {"sku": r["sku"], "name": r["name"], "qty": r["qty"],
                      "box": r["box"], "expire": r.get("expire", "")}
                g["rows"].append(nr)
                by_sku[r["sku"]] = nr
        parsed_files.append({"name": up.filename,
                             "kind": f"거래명세서({disp}, {len(d['rows'])}품목)"})

    if not centers:
        raise HTTPException(
            400, "거래명세서를 인식하지 못했습니다. 쿠팡 거래명세서 PDF를 포함해주세요."
        )

    # 센터 자연 정렬
    ordered = sorted(centers.items(), key=lambda kv: _center_sortkey(kv[0]))

    groups = []
    text_blocks = ["[출고 수량 상세]", ""]
    for disp, g in ordered:
        items = [{"name": r["name"], "qty": r["qty"], "box": r["box"],
                  "expire": r.get("expire", "")} for r in g["rows"]]
        groups.append({
            "center": disp,
            "date": g.get("date", ""),
            "items": items,
            "totalQty": sum(r["qty"] for r in g["rows"]),
            "totalBox": sum(r["box"] for r in g["rows"]),
        })
        text_blocks.append(f"쿠팡 {disp} 센터")
        text_blocks.append("")
        for r in g["rows"]:
            line = f"{r['name']}: {r['qty']:,}개 ({r['box']}박스)"
            if r.get("expire"):
                line += f" / 유통기한 {r['expire']} 이후"
            text_blocks.append(line)
            text_blocks.append("")

    return {
        "text": "\n".join(text_blocks).rstrip() + "\n",
        "groups": groups,
        "summary": {
            "centers": len(groups),
            "totalQty": sum(g["totalQty"] for g in groups),
            "totalBox": sum(g["totalBox"] for g in groups),
        },
        "files": parsed_files,
        "parse_errors": errors,
    }


# ── 마켓컬리: 라벨지(PPTX) 파서 + 거래명세서 대조 ─────────────────────────────
def parse_kurly_label_pptx(data: bytes) -> dict:
    """입고 라벨지 PPTX: 슬라이드 1장 = 박스 라벨 1장. 상품코드별로 라벨 수를 집계한다."""
    z = zipfile.ZipFile(io.BytesIO(data))
    slide_files = [n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)]
    by_code: dict[str, dict] = {}
    order_codes = set()
    total = 0
    for fn in slide_files:
        xml = z.read(fn).decode("utf-8", "ignore")
        texts = [html.unescape(t) for t in re.findall(r"<a:t>(.*?)</a:t>", xml, re.S)]
        joined = " ".join(texts)
        mcode = re.search(r"(M\d{6,})", joined)
        if not mcode:
            continue
        total += 1
        code = mcode.group(1)
        moc = re.search(r"발주코드\s*([A-Za-z0-9_]+)", joined)
        if moc:
            order_codes.add(moc.group(1))
        mname = re.search(r"상품명\s*(.+?)\s*상품코드", joined)
        mpb = re.search(r"입수량\s*\(\s*(\d+)\s*\).*?총\s*입고수량\s*\(\s*(\d+)\s*\)", joined)
        mbt = re.search(r"전체\s*박스\s*수\s*\(\s*(\d+)\s*\)", joined)
        rec = by_code.setdefault(code, {"name": "", "perBox": None, "total": None, "boxTotal": None, "count": 0})
        rec["count"] += 1
        if mname and not rec["name"]:
            rec["name"] = norm_product(mname.group(1))
        if mpb:
            rec["perBox"], rec["total"] = int(mpb.group(1)), int(mpb.group(2))
        if mbt:
            rec["boxTotal"] = int(mbt.group(1))
    return {
        "orderCode": sorted(order_codes)[0] if order_codes else "",
        "orderCodes": sorted(order_codes),
        "byCode": by_code,
        "total_slides": total,
    }


def parse_kurly_request(data: bytes) -> dict:
    """컬리 출고요청 엑셀(택배 라벨용). 주문번호=발주코드, 한 행=한 박스, 수량=박스당입수.

    반환: { 발주코드: {center, byName: {pkey(상품명): {name, box, total, perBox}}} }
    """
    xl = pd.ExcelFile(io.BytesIO(data))
    df = xl.parse(xl.sheet_names[0])
    by_order: dict[str, dict] = {}
    for _, r in df.iterrows():
        oc = str(r.get("주문번호") or "").strip()
        if not oc or oc.lower() == "nan":
            continue
        name = norm_product(r.get("상품명"))
        if not name:
            continue
        qty = _to_int(r.get("수량"))
        g = by_order.setdefault(oc, {"center": norm_center(r.get("수취인")), "byName": {}})
        rec = g["byName"].setdefault(pkey(name), {"name": name, "box": 0, "total": 0, "perboxes": set()})
        rec["box"] += 1            # 한 행 = 한 박스(=라벨 1장)
        rec["total"] += qty
        if qty:
            rec["perboxes"].add(qty)
    for g in by_order.values():
        for rec in g["byName"].values():
            pbs = rec.pop("perboxes")
            rec["perBox"] = sorted(pbs)[-1] if pbs else 0
    return by_order


def build_kurly_group(order_code: str, stmt: dict | None, labels: dict | None,
                      req: dict | None = None) -> dict:
    checks = []
    statuses = []

    present = []
    if stmt and stmt.get("items"):
        present.append("거래명세서")
    if labels and labels.get("total_slides"):
        present.append("라벨지")
    if req and req.get("byName"):
        present.append("출고요청")
    missing_required = [k for k in ("거래명세서", "라벨지") if k not in present]
    if missing_required:
        checks.append({"label": "서류 구비", "status": "error", "detail": f"누락: {', '.join(missing_required)}"})
        statuses.append("error")
    else:
        detail = " + ".join(present)
        if "출고요청" not in present:
            detail += " (출고요청 없음)"
        checks.append({"label": "서류 구비", "status": "ok", "detail": detail})

    # 발주코드 일치
    ocs = sorted({x for x in ([stmt.get("orderCode")] if stmt else []) + (labels.get("orderCodes", []) if labels else []) if x})
    if len(ocs) <= 1:
        checks.append({"label": "발주코드 일치", "status": "ok", "detail": order_code or (ocs[0] if ocs else "-")})
    else:
        checks.append({"label": "발주코드 일치", "status": "error", "detail": " ≠ ".join(ocs)})
        statuses.append("error")

    stmt_items = {it["code"]: it for it in (stmt["items"] if stmt else [])}
    label_codes = labels.get("byCode", {}) if labels else {}
    req_by_name = req.get("byName", {}) if req else {}
    has_req = bool(req_by_name)
    req_used = set()

    items = []
    tot = {"reqBox": 0, "stmtBox": 0, "labelCount": 0, "reqTotal": 0, "stmtTotal": 0, "labelTotal": 0}
    sku_error = False
    for code in sorted(set(stmt_items) | set(label_codes)):
        si = stmt_items.get(code)
        li = label_codes.get(code)
        name = (si["name"] if si else "") or (li["name"] if li else "") or "(이름 미상)"
        ri = req_by_name.get(pkey(name))
        if ri:
            req_used.add(pkey(name))

        req_box = ri["box"] if ri else 0
        req_total = ri["total"] if ri else 0
        stmt_box = si["boxCount"] if si else 0
        label_count = li["count"] if li else 0
        stmt_total = si["total"] if si else 0
        label_total = li["total"] if (li and li.get("total") is not None) else 0
        stmt_perbox = si["perBox"] if si else 0
        label_perbox = li["perBox"] if (li and li.get("perBox") is not None) else 0

        notes = []
        st = "ok"
        if not si:
            st = "error"; notes.append("거래명세서에 없음")
        elif not li:
            st = "error"; notes.append("라벨 없음")
        else:
            if label_count != stmt_box:
                st = "error"; notes.append(f"라벨 {label_count}장 ≠ 박스수 {stmt_box}")
            if stmt_total != label_total:
                st = "error"; notes.append(f"총수량 명세 {stmt_total} ≠ 라벨 {label_total}")
            if li.get("boxTotal") is not None and li["boxTotal"] != label_count:
                if st == "ok":
                    st = "warn"
                notes.append(f"라벨표기 박스수 {li['boxTotal']} ≠ 실제 {label_count}장")
            if stmt_perbox != label_perbox:
                if st == "ok":
                    st = "warn"
                notes.append(f"입수 명세 {stmt_perbox} ≠ 라벨 {label_perbox}")
        # 출고요청 ↔ 거래명세서(확정) 비교 (요청 정보 있을 때)
        if ri and si:
            if req_box != stmt_box or req_total != stmt_total:
                if st == "ok":
                    st = "warn"
                notes.append(f"요청 {req_box}박스/{req_total}개 ≠ 명세 {stmt_box}박스/{stmt_total}개")
        elif has_req and si and not ri:
            if st == "ok":
                st = "warn"
            notes.append("출고요청에 없음")

        if st == "error":
            sku_error = True
        items.append({
            "code": code, "name": name,
            "reqBox": req_box, "stmtBox": stmt_box, "labelCount": label_count,
            "reqTotal": req_total, "stmtTotal": stmt_total, "labelTotal": label_total,
            "stmtPerBox": stmt_perbox, "labelPerBox": label_perbox,
            "status": st, "note": ", ".join(notes),
        })
        tot["reqBox"] += req_box
        tot["stmtBox"] += stmt_box
        tot["labelCount"] += label_count
        tot["reqTotal"] += req_total
        tot["stmtTotal"] += stmt_total
        tot["labelTotal"] += label_total

    # 출고요청에만 있는 품목(명세/라벨에 매칭 안 됨)
    for pk, ri in req_by_name.items():
        if pk in req_used:
            continue
        sku_error = True
        items.append({
            "code": "", "name": ri["name"],
            "reqBox": ri["box"], "stmtBox": 0, "labelCount": 0,
            "reqTotal": ri["total"], "stmtTotal": 0, "labelTotal": 0,
            "stmtPerBox": 0, "labelPerBox": 0,
            "status": "error", "note": "출고요청에만 있음 (명세/라벨 없음)",
        })
        tot["reqBox"] += ri["box"]
        tot["reqTotal"] += ri["total"]

    if sku_error:
        statuses.append("error")

    # 총 라벨 수 = 총 박스 수 (+ 요청 포함)
    box_vals = {tot["stmtBox"], tot["labelCount"]}
    if has_req:
        box_vals.add(tot["reqBox"])
    if len(box_vals) == 1:
        parts = []
        if has_req:
            parts.append(f"요청 {tot['reqBox']}")
        parts += [f"명세 {tot['stmtBox']}", f"라벨 {tot['labelCount']}"]
        checks.append({"label": "총 박스/라벨 수", "status": "ok", "detail": " = ".join(parts) + "박스"})
    else:
        st = "error" if tot["stmtBox"] != tot["labelCount"] else "warn"
        parts = []
        if has_req:
            parts.append(f"요청 {tot['reqBox']}")
        parts += [f"명세 {tot['stmtBox']}", f"라벨 {tot['labelCount']}"]
        checks.append({"label": "총 박스/라벨 수", "status": st, "detail": " / ".join(parts)})
        statuses.append(st)

    status = _worst(*statuses) if statuses else "ok"
    return {
        "orderCode": order_code,
        "center": (stmt.get("center") if stmt else "") or (req.get("center") if req else "") or "",
        "date": (stmt.get("date") if stmt else "") or "",
        "supplier": (stmt.get("supplier") if stmt else "") or "",
        "status": status,
        "present": present,
        "missing": missing_required,
        "checks": checks,
        "items": items,
        "totals": tot,
        "hasReq": has_req,
    }


@router.post("/kurly")
async def review_kurly(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(400, "파일을 하나 이상 업로드해주세요.")

    stmts: dict[str, dict] = {}
    labels: dict[str, dict] = {}
    reqs: dict[str, dict] = {}
    parsed_files = []
    errors = []

    for up in files:
        if not up.filename:
            continue
        data = await up.read()
        name = up.filename.lower()
        try:
            if name.endswith((".xlsx", ".xls")):
                by_order = parse_kurly_request(data)
                reqs.update(by_order)
                parsed_files.append({"name": up.filename, "kind": f"출고요청({len(by_order)}발주)"})
            elif name.endswith((".pptx", ".ppt")) or "라벨" in up.filename:
                d = parse_kurly_label_pptx(data)
                oc = d["orderCode"] or up.filename
                labels[oc] = d
                parsed_files.append({"name": up.filename, "kind": "라벨지"})
            elif name.endswith(".pdf"):
                d = kurly_label.parse_kurly_statement(data)
                oc = d["orderCode"] or up.filename
                stmts[oc] = d
                parsed_files.append({"name": up.filename, "kind": "거래명세서"})
            else:
                parsed_files.append({"name": up.filename, "kind": "기타"})
        except Exception as e:  # noqa: BLE001
            errors.append(f"{up.filename}: {e}")

    ocs = sorted(set(stmts) | set(labels) | set(reqs))
    if not ocs:
        raise HTTPException(
            400,
            "마켓컬리 거래명세서(PDF)·입고 라벨지(PPTX)·출고요청(Excel)을 인식하지 못했습니다.",
        )

    groups = [build_kurly_group(oc, stmts.get(oc), labels.get(oc), reqs.get(oc)) for oc in ocs]
    summary = {
        "total": len(groups),
        "ok": sum(1 for g in groups if g["status"] == "ok"),
        "warn": sum(1 for g in groups if g["status"] == "warn"),
        "error": sum(1 for g in groups if g["status"] == "error"),
    }
    return {
        "groups": groups,
        "summary": summary,
        "files": parsed_files,
        "parse_errors": errors,
    }
