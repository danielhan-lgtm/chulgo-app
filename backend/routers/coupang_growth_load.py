# -*- coding: utf-8 -*-
"""쿠팡 그로스(로켓그로스) 팔레트 적재리스트 PPT 생성기.

동봉문서(거래명세서, 한글 쿠팡 제출용) + 부착문서(로켓그로스 팔레트 라벨 + 적재리스트 양식)를
업로드하면, 부착문서 마지막 장의 '쿠팡 팔레트 적재리스트' 양식에 내용을 채운 PPT를 생성한다.

- 동봉문서 → 업체명/업체코드 + 발주/물류센터/도착예정일 + 상품표(SKU·상품명·수량·소비기한/제조일자)
- 부착문서 → 요청ID(밀크런ID) · 물류센터 · 도착예정일 · 팔레트번호 · 박스바코드
PPT는 python-pptx로 양식(업체정보/입고예약정보/상품표)을 직접 그려 생성한다. 외부 API 미사용.
"""
import io
import os
import re
import sys

import pdfplumber
from fastapi import APIRouter, File, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import load_settings as LS

router = APIRouter()


def _clean(s) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()


def _to_int(v) -> int:
    """'4,680'→4680, 4680.0/'4680.0'→4680, NaN/빈값→0 (마침표를 자릿수로 붙이는 오류 방지)."""
    if v is None:
        return 0
    if isinstance(v, (int, float)):
        try:
            return int(round(v)) if v == v else 0  # NaN != NaN
        except Exception:
            return 0
    s = str(v).replace(",", "").strip()
    if not s:
        return 0
    try:
        return int(round(float(s)))
    except Exception:
        pass
    try:
        return int(re.sub(r"[^\d\-]", "", s) or 0)
    except Exception:
        return 0


# 상품별 박스당 입수 — 저장된 설정(load_settings.json)의 규칙/기본값 사용 (기본 132개/박스)


def _fmt_date(s: str) -> str:
    """'20260707' 또는 '2026-07-07' → '2026-07-07'."""
    s = str(s or "").strip()
    m = re.search(r"(\d{4})[-.]?(\d{2})[-.]?(\d{2})", s)
    return f"{m.group(1)}-{m.group(2)}-{m.group(3)}" if m else s


# ── 동봉문서(거래명세서) 파서 ───────────────────────────────────────────────
def parse_동봉문서(data: bytes) -> dict:
    supplier = supplier_code = po = center = date = ""
    products: list[dict] = []
    seen_sku: set[str] = set()

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for pg in pdf.pages:
            for tbl in pg.extract_tables():
                for row in tbl:
                    if not row:
                        continue
                    cells = [(_clean(c)) for c in row]
                    key = cells[0]
                    val = cells[1] if len(cells) > 1 else ""
                    # 거래처/발주 정보 (key-value 표)
                    if key in ("업체명",) and val and not supplier:
                        supplier = val
                    elif key in ("업체번호",) and val and not supplier_code:
                        supplier_code = val
                    elif key == "발주번호" and val and not po:
                        po = val
                    elif key in ("납품 센터", "납품센터") and val and not center:
                        center = val
                    elif "도착예정일" in key and val and not date:
                        date = _fmt_date(val)
                    # 상품행: col0=숫자 No, col1=상품번호(숫자)
                    if len(cells) >= 7 and cells[0].isdigit() and cells[1].isdigit():
                        sku = cells[1]
                        if sku in seen_sku:
                            continue
                        seen_sku.add(sku)
                        # col2: 상품명/옵션 줄 ... 바코드줄 / Box바코드줄 (줄바꿈 구분)
                        # 원본(줄바꿈 유지)에서 바코드 줄 전까지만 상품명으로 사용
                        raw_name = str(row[2] or "")
                        name_parts = []
                        for line in raw_name.split("\n"):
                            line = line.strip()
                            if not line or line == "-":
                                continue
                            if re.fullmatch(r"[A-Z]?\d{9,}", line):  # 상품/Box 바코드 줄
                                break
                            name_parts.append(line)
                        name = _clean(" ".join(name_parts))
                        qty = _to_int(cells[4]) or _to_int(cells[3])  # 확정수량 우선
                        # 제조(수입일자)/소비기한 (col6 'YYYYMMDD/YYYYMMDD')
                        dts = re.findall(r"\d{8}|\d{4}-\d{2}-\d{2}", cells[6] if len(cells) > 6 else "")
                        made = _fmt_date(dts[0]) if dts else ""
                        expire = _fmt_date(dts[1]) if len(dts) > 1 else ""
                        products.append({
                            "sku": sku, "name": name, "qty": qty,
                            "made": made, "expire": expire,
                        })
    return {
        "supplier": supplier, "supplier_code": supplier_code, "po": po,
        "center": center, "date": date, "products": products,
    }


# ── 부착문서(로켓그로스) 파서 ───────────────────────────────────────────────
def parse_부착문서(data: bytes) -> dict:
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        txt = "\n".join((pg.extract_text() or "") for pg in pdf.pages)

    out = {"milkrun_id": "", "center": "", "date": "", "pallet": "", "pallet_total": 1,
           "box_barcode": "", "supplier": "", "supplier_code": ""}

    # 물류센터: 'XXX(NN)' 바로 뒤 '[로켓그로...' (줄바꿈으로 쪼개진 경우 대비, 센터/팔레트 분리 파싱)
    m = re.search(r"([^\s\[\]]+\([0-9A-Za-z]+\))\s*\[?\s*로켓그로", txt)
    if m:
        out["center"] = m.group(1).strip()
    # 팔레트 번호: '팔레트 3-1' (어디에 있든) — 첫 숫자가 총 팔레트 수
    m = re.search(r"팔레트\s*(\d+\s*-\s*\d+)", txt)
    if m:
        out["pallet"] = re.sub(r"\s+", "", m.group(1))
        out["pallet_total"] = max(1, int(re.findall(r"\d+", m.group(1))[0]))
    # '밀크런ID ... 도착예정일' 다음 줄 '10540793 2026-07-07'
    m = re.search(r"밀크런ID[^\n]*\n\s*(\d+)\s+([\d\-]+)", txt)
    if m:
        out["milkrun_id"] = m.group(1)
        out["date"] = _fmt_date(m.group(2))
    else:
        m = re.search(r"밀크런ID\s*[:：]?\s*(\d+)", txt)
        if m:
            out["milkrun_id"] = m.group(1)
        m = re.search(r"도착예정일\s*[:：]?\s*([\d\-]+)", txt)
        if m:
            out["date"] = _fmt_date(m.group(1))
    m = re.search(r"(MRN\w+)", txt)
    if m:
        out["box_barcode"] = m.group(1)
    # 업체명/업체코드 (적재리스트 양식 또는 라벨)
    m = re.search(r"업체코드\s*([A-Z]\d+)", txt)
    if m:
        out["supplier_code"] = m.group(1)
    return out


def _sniff(filename: str, data: bytes) -> str:
    fn = filename or ""
    if "동봉" in fn or "명세" in fn:
        return "동봉"
    if "부착" in fn:
        return "부착"
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            head = (pdf.pages[0].extract_text() or "") if pdf.pages else ""
    except Exception:
        return "기타"
    if "팔레트 부착" in head or "로켓그로스" in head or "밀크런ID" in head:
        return "부착"
    if "거래 명세서" in head or "발주번호" in head or "업체번호" in head:
        return "동봉"
    return "기타"


# ── 파싱 엔드포인트 ────────────────────────────────────────────────────────
@router.post("/parse")
async def parse_docs(files: list[UploadFile] = File(...)):
    if not files:
        raise HTTPException(400, "동봉문서/부착문서 PDF를 업로드해주세요.")
    동봉: dict | None = None
    부착: dict | None = None
    errors: list[str] = []
    for up in files:
        if not up.filename:
            continue
        data = await up.read()
        kind = _sniff(up.filename, data)
        try:
            if kind == "부착":
                부착 = parse_부착문서(data)
            elif kind == "동봉":
                d = parse_동봉문서(data)
                if 동봉 is None:
                    동봉 = d
                else:
                    동봉["products"].extend(d["products"])
            else:
                errors.append(f"{up.filename}: 종류 인식 실패")
        except Exception as e:  # noqa: BLE001
            errors.append(f"{up.filename}: {e}")

    if 동봉 is None and 부착 is None:
        raise HTTPException(400, "PDF를 인식하지 못했습니다. 동봉문서/부착문서를 확인해주세요. "
                            + (" / ".join(errors) if errors else ""))
    동봉 = 동봉 or {"supplier": "", "supplier_code": "", "po": "", "center": "", "date": "", "products": []}
    부착 = 부착 or {}

    box_barcode = 부착.get("box_barcode", "")
    settings = LS.get_settings("growth")
    rows = []
    for i, p in enumerate(동봉.get("products", []), start=1):
        rows.append({
            "no": i, "sku": p["sku"], "name": p["name"],
            "box_no": box_barcode,           # 박스 바코드(참고용)
            "box": LS.guess_box(settings, p["sku"], p["name"], p["qty"]),  # 박스 수(설정 기반, 수정가능)
            "qty": p["qty"],
            "expire": p.get("expire", ""), "made": p.get("made", ""),
        })

    return {
        "supplier": 동봉.get("supplier", ""),
        "supplier_code": 동봉.get("supplier_code", "") or 부착.get("supplier_code", ""),
        "request_id": 부착.get("milkrun_id", "") or 동봉.get("po", ""),
        "center": 부착.get("center", "") or 동봉.get("center", ""),
        "date": 부착.get("date", "") or 동봉.get("date", ""),
        "pallet": 부착.get("pallet", "") or "1-1",
        "pallet_total": 부착.get("pallet_total", 1) or 1,
        "box_barcode": box_barcode,
        "rows": rows,
        "parse_errors": errors,
        "has_부착": bool(부착.get("milkrun_id") or 부착.get("center")),
        "settings": settings,
    }


# ── 적재 설정 (입수 규칙) ──────────────────────────────────────────────────
class GLoadRule(BaseModel):
    match: str = ""       # 상품번호(전체 일치) 또는 상품명 키워드(포함)
    per_box: int = 1      # 1박스당 수량


class GrowthLoadSettings(BaseModel):
    pallet_cap: int = 112
    pallet_total_mode: str = "auto"   # auto | attach | fixed
    pallet_total_fixed: int = 1
    default_per_box: int = 132
    bundle_is_box: bool = False
    rules: list[GLoadRule] = []


@router.get("/settings")
def get_load_settings():
    return LS.get_settings("growth")


@router.post("/settings")
def save_load_settings(body: GrowthLoadSettings):
    return LS.save_settings("growth", body.model_dump())


# ── PPT 생성 ───────────────────────────────────────────────────────────────
class GRow(BaseModel):
    no: int | None = None
    sku: str = ""
    name: str = ""
    box_no: str = ""
    box: int = 0
    qty: int = 0
    expire: str = ""
    made: str = ""
    pallet: int = 1          # 이 품목이 실릴 팔레트(=슬라이드) 번호


class GPayload(BaseModel):
    supplier: str = ""
    supplier_code: str = ""
    request_id: str = ""
    center: str = ""
    date: str = ""
    pallet: str = "1-1"     # (호환용) 단일 팔레트일 때의 표기
    pallet_total: int = 1   # 총 팔레트 수 = 생성할 슬라이드 장수
    total_box: int = 0
    rows: list[GRow] = []


def _plain_grid(tbl):
    """python-pptx 기본 표 스타일(색 밴딩) 제거 → 'No Style, Table Grid'(흰 셀+검정 테두리)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    tblPr = tbl._tbl.tblPr
    for a in ("firstRow", "bandRow", "firstCol", "lastCol", "lastRow", "bandCol"):
        tblPr.set(a, "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is None:
        sid = etree.SubElement(tblPr, qn("a:tableStyleId"))
    sid.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # No Style, Table Grid


def _draw_slide(prs, p: GPayload, pallet_label: str, rows: list, total_box: int):
    """팔레트 1개(=슬라이드 1장)에 업체정보/입고예약정보/상품표를 그린다."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    L, W = Inches(0.4), Inches(7.47)

    def textbox(top, text, size, bold=False, align=PP_ALIGN.LEFT, height=0.3, left=L, width=W):
        tb = slide.shapes.add_textbox(left, top, width, Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = Pt(0)
        para = tf.paragraphs[0]; para.alignment = align
        run = para.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = BLACK
        run.font.name = "맑은 고딕"
        return tb

    def cell_style(cell, text, size=9, bold=False, align=PP_ALIGN.CENTER):
        cell.text = "" if text is None else str(text)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Pt(3)
        cell.margin_top = cell.margin_bottom = Pt(1)
        for para in cell.text_frame.paragraphs:
            para.alignment = align
            for r in para.runs:
                r.font.size = Pt(size); r.font.bold = bold
                r.font.color.rgb = BLACK; r.font.name = "맑은 고딕"

    # 제목
    textbox(Inches(0.3), "쿠팡 팔레트 적재리스트 (각 팔레트 부착 필수)", 15, True, PP_ALIGN.CENTER, 0.4)
    textbox(Inches(0.72), "※ 팔레트의 높이는 1,700mm를 초과할 수 없습니다 ※", 9, False, PP_ALIGN.CENTER, 0.25)

    def kv_table(top, title, pairs):
        textbox(top, title, 11, True, PP_ALIGN.LEFT, 0.26)
        rows_n = len(pairs)
        gt = slide.shapes.add_table(rows_n, 2, L, top + Inches(0.3), W, Inches(0.28 * rows_n)).table
        _plain_grid(gt)
        gt.columns[0].width = Inches(2.2)
        gt.columns[1].width = W - Inches(2.2)
        for i, (k, v) in enumerate(pairs):
            gt.rows[i].height = Inches(0.28)
            cell_style(gt.cell(i, 0), k, 10, True, PP_ALIGN.CENTER)
            cell_style(gt.cell(i, 1), v, 10, False, PP_ALIGN.LEFT)
        return top + Inches(0.3) + Inches(0.28 * rows_n)

    y = kv_table(Inches(1.05), "1 ) 업체 정보",
                 [("업체명", p.supplier), ("업체코드", p.supplier_code)])
    y = kv_table(y + Inches(0.2), "2 ) 입고 예약 정보",
                 [("요청 ID", p.request_id), ("물류센터", p.center),
                  ("물류센터 도착예정일", p.date), ("팔레트 번호", pallet_label),
                  ("총 박스", str(total_box))])

    # 3) 상품 정보 — 최소 12행 고정(양식과 동일)
    textbox(y + Inches(0.2), "3 ) 상품 정보", 11, True, PP_ALIGN.LEFT, 0.26)
    headers = ["No.", "SKU ID", "물류입고용 상품명 / 옵션명", "박스 번호", "상품 수량", "소비기한/제조일자"]
    body = max(12, len(rows))
    n = body + 1
    tbl = slide.shapes.add_table(n, len(headers), L, y + Inches(0.5), W, Inches(0.3 * n)).table
    _plain_grid(tbl)
    for ci, wd in enumerate([Inches(0.5), Inches(1.15), Inches(3.02), Inches(0.9), Inches(0.9), Inches(1.0)]):
        tbl.columns[ci].width = wd
    for ci, h in enumerate(headers):
        cell_style(tbl.cell(0, ci), h, 9, True, PP_ALIGN.CENTER)
    for i in range(1, n):
        tbl.rows[i].height = Inches(0.28)
        row = rows[i - 1] if i - 1 < len(rows) else None
        if row is None:
            cell_style(tbl.cell(i, 0), i, 9, False, PP_ALIGN.CENTER)
            for ci in range(1, len(headers)):
                cell_style(tbl.cell(i, ci), "", 9)
            continue
        exp = row.expire + (f" / {row.made}" if row.made else "")
        vals = [i, row.sku, row.name, (row.box or ""), row.qty, exp]
        for ci, val in enumerate(vals):
            cell_style(tbl.cell(i, ci), val, 9, False, PP_ALIGN.LEFT if ci == 2 else PP_ALIGN.CENTER)


def build_pptx(p: GPayload) -> io.BytesIO:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(8.27)     # A4 세로 (595pt)
    prs.slide_height = Inches(11.69)   # (842pt)

    rows = p.rows or []
    # 총 팔레트 수 = 지정값과 실제 사용된 최대 팔레트 번호 중 큰 값(최소 1)
    max_used = max((int(r.pallet or 1) for r in rows), default=1)
    total = max(int(p.pallet_total or 1), max_used, 1)

    # 팔레트 번호별 그룹핑(범위를 벗어난 값은 1번으로)
    groups: dict[int, list] = {i: [] for i in range(1, total + 1)}
    for r in rows:
        pal = int(r.pallet or 1)
        if pal < 1 or pal > total:
            pal = 1
        groups[pal].append(r)

    for idx in range(1, total + 1):
        grp = groups[idx]
        box_total = sum(_to_int(r.box) for r in grp)
        if total == 1:
            # 단일 팔레트: 부착문서 표기(예: '1-1')와 화면의 총박스 입력값 유지
            label = p.pallet or "1-1"
            box_total = p.total_box or box_total
        else:
            label = f"{total}-{idx}"
        _draw_slide(prs, p, label, grp, box_total)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf


@router.post("/generate")
def generate(payload: GPayload):
    from urllib.parse import quote
    buf = build_pptx(payload)
    name = f"{payload.center or '쿠팡그로스'}_적재리스트.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=growth_load.pptx; filename*=UTF-8''{quote(name)}"},
    )
